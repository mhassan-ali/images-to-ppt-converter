from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# ================== YOUR INFO ==================
YOUR_NAME    = "M.Hassan.Ali"
ROLL_NO      = "2024F-BCE-149"
COURSE       = "CE-207T Computer Organization & Architecture"
TEACHER      = "M. Asim Hasan"
IMAGE_FOLDER = "images"        # <-- your images are inside images/ folder
OUTPUT_FILE  = "USB_Presentation_M_Hassan_Ali.pptx"
# ===============================================

def find_image(folder, number):
    """Try to find image with different naming formats"""
    formats = [
        f"{number:02d}",        # 01, 02, 03
        f"{number}",            # 1, 2, 3
        f"slide-{number:02d}",  # slide-01
        f"slide_{number:02d}",  # slide_01
        f"Slide{number:02d}",   # Slide01
        f"Slide{number}",       # Slide1
    ]
    extensions = ['.jpg', '.jpeg', '.png', '.JPG', '.JPEG', '.PNG']

    for name in formats:
        for ext in extensions:
            path = os.path.join(folder, f"{name}{ext}")
            if os.path.exists(path):
                return path
    return None


def create_ppt():
    print("\n" + "="*50)
    print("  USB PRESENTATION GENERATOR")
    print("="*50)

    # Check if images folder exists
    if not os.path.exists(IMAGE_FOLDER):
        print(f"\n❌ ERROR: '{IMAGE_FOLDER}' folder not found!")
        print(f"   Current directory: {os.getcwd()}")
        print(f"   Make sure your images folder is here:")
        print(f"   {os.path.abspath(IMAGE_FOLDER)}")
        return

    # Show what files are found in images folder
    print(f"\n📁 Looking in folder: {os.path.abspath(IMAGE_FOLDER)}")
    files_in_folder = os.listdir(IMAGE_FOLDER)
    print(f"📷 Files found: {files_in_folder}")

    # Create presentation (16:9)
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides_added = 0

    for i in range(1, 16):
        image_path = find_image(IMAGE_FOLDER, i)

        if not image_path:
            print(f"\n⚠️  Slide {i:02d} → Image NOT found! Skipping...")
            continue

        print(f"✅ Slide {i:02d} → {image_path}")

        # Add blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Dark background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(15, 15, 26)

        # ---- Top blue line ----
        top_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(13.33), Inches(0.05)
        )
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = RGBColor(0, 212, 255)
        top_line.line.fill.background()

        # ---- Slide number (top left) ----
        num_box = slide.shapes.add_textbox(
            Inches(0.2), Inches(0.08),
            Inches(0.6), Inches(0.4)
        )
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = f"{i:02d}"
        num_p.font.size = Pt(13)
        num_p.font.bold = True
        num_p.font.color.rgb = RGBColor(0, 212, 255)
        num_p.font.name = "Consolas"

        # ---- Slide counter top right (1/15) ----
        counter_box = slide.shapes.add_textbox(
            Inches(12.3), Inches(0.08),
            Inches(0.8), Inches(0.4)
        )
        counter_tf = counter_box.text_frame
        counter_p = counter_tf.paragraphs[0]
        counter_p.text = f"{i}/15"
        counter_p.font.size = Pt(11)
        counter_p.font.color.rgb = RGBColor(100, 100, 120)
        counter_p.font.name = "Consolas"
        counter_p.alignment = PP_ALIGN.RIGHT

        # ---- Main Image ----
        slide.shapes.add_picture(
            image_path,
            Inches(0.15), Inches(0.55),
            width=Inches(13.03)
        )

        # ---- Bottom green line ----
        bottom_line = slide.shapes.add_shape(
            1,
            Inches(0), Inches(7.45),
            Inches(13.33), Inches(0.05)
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = RGBColor(0, 255, 136)
        bottom_line.line.fill.background()

        # ---- Footer: Name | Roll | Course ----
        footer_box = slide.shapes.add_textbox(
            Inches(0.2), Inches(7.1),
            Inches(13), Inches(0.35)
        )
        footer_tf = footer_box.text_frame
        footer_p = footer_tf.paragraphs[0]
        footer_p.text = f"{YOUR_NAME}  |  {ROLL_NO}  |  {COURSE}  |  {TEACHER}"
        footer_p.font.size = Pt(9)
        footer_p.font.color.rgb = RGBColor(80, 80, 100)
        footer_p.font.name = "Calibri"
        footer_p.alignment = PP_ALIGN.CENTER

        slides_added += 1

    # Save
    if slides_added == 0:
        print("\n❌ No slides were created!")
        print("   Check that your images are named: 01.jpg, 02.jpg ... 15.jpg")
        print(f"   And placed inside the '{IMAGE_FOLDER}' folder")
        return

    prs.save(OUTPUT_FILE)

    print(f"\n{'='*50}")
    print(f"  ✅ SUCCESS!")
    print(f"  📊 Slides created : {slides_added}/15")
    print(f"  📁 Saved as       : {OUTPUT_FILE}")
    print(f"  👤 Name           : {YOUR_NAME}")
    print(f"  🎓 Roll No        : {ROLL_NO}")
    print(f"{'='*50}\n")


# Run
create_ppt()